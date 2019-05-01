#!/usr/bin/env python
# coding: utf-8

# # 为PPT配音
# 

# ## 初始化

# In[1]:


import sys, os , subprocess, time, platform, json
if platform.system()=="Darwin":
    try:
        from  AppKit import NSSpeechSynthesizer
        import Foundation
    except:
        pass
from pptx import Presentation
from pptx.util import Inches
# from xunfei_tts import init_xf_tts, xf_save_tts
from google_tts import Speech as gSpeech
from xunfei_tts import Speech as xf_Speech
from pydub import AudioSegment


# In[2]:


def read_pptx(ppt_filename):
    prs = Presentation(ppt_filename)
    return prs

def init_tts(tts_engine):
    if tts_engine=="nsss":
        nssp = NSSpeechSynthesizer
        ve = nssp.alloc().init()
    elif tts_engine=="espeak":
        ve = 0
    elif tts_engine=="sapi5":
        ve = 0
    elif tts_engine=="xunfei":
#         ve = init_xf_tts()
        with open('API_setup.txt') as json_file:  
            api = json.load(json_file)
            ve = xf_Speech(api, voice_name="x_xiaofeng")
    elif tts_engine=="google":
        ve = 0
    else:
        ve = 0
    return ve


# ## 读取ppt中每页的注释

# In[3]:


def get_notes_text(slide):
    notes_slide = slide.notes_slide
    note = notes_slide.notes_text_frame.text
    return note


# ## 将每页注释用tts转换成音频文件

# In[4]:


def save_tts(ve, TEXT, filename, tts_engine):
    if tts_engine=="nsss":
        filename=filename+".aiff"
        url = Foundation.NSURL.fileURLWithPath_(filename)
        s=ve.startSpeakingString_toURL_(TEXT,url)    
        if not(s):
            print("TTS failed") 
    elif tts_engine=="espeak":
        # 尚未测试
        filename=filename+".wav"
        subprocess.call(["espeak",TEXT,"-w", filename])
    elif tts_engine=="sapi5":
        # 看起来很复杂的样子, 参考: https://github.com/nateshmbhat/pyttsx3/issues/7
        pass
    elif tts_engine=="xunfei":
        filename=filename+".wav"
#         xf_save_tts(ve, TEXT, filename)
        ve.save(TEXT,filename)
    elif tts_engine=="google":
        filename=filename+".mp3"
        ve = gSpeech(TEXT, "zh-cn","en")
        ve.save(filename)
    return filename


# In[5]:


def save_notes_voice(ve, text, page_number, tts_engine):
    voice_filename= "temp_tts_{:3d}".format(page_number)
    voice_filename= save_tts(ve, text,voice_filename, tts_engine )
    return voice_filename


# 从tts引擎传回的音频文件可能存在编码问题。于是pydub.AudioSegment再重新灌注一边好了。
# 

# In[ ]:


def reclean_voice(voice_filename):
    _, file_extension = os.path.splitext(voice_filename)
    file_extension=file_extension.lstrip(".")
    (AudioSegment.from_file(voice_filename, 
                           format=file_extension)
                .export(voice_filename, 
                        format=file_extension,
                        bitrate="32k"))
    new_voice_filename=voice_filename

# # ffmpeg -i input.wav -vn -ar 44100 -ac 2 -b:a 192k output.mp3
#     file_name, file_extension = os.path.splitext(voice_filename)
#     new_voice_filename=file_name+"_new"+'.aif'
#     subprocess.call(['ffmpeg',
#                      '-i', voice_filename,
# #                      '-vn','-ar','44100','-ac','1','-b:a','16k',
# #                      '-acodec', 'libmp3lame',
#                      new_voice_filename
#                     ])
    return new_voice_filename
    


# ## 将每个音频文件插入到ppt页面中

# In[6]:


def insert_voice(voice_filename, slide):
    left = top = Inches(0.0)
    width = height = Inches(1.0)

    shapes = slide.shapes
    movie = shapes.add_movie(voice_filename, 
                                 left , top , width , height, 
                                 poster_frame_image=None, 
                                 mime_type='video/unknown')


# ## 清理掉临时文件

# In[7]:


def clean_temp(voice_filename):
    os.remove(voice_filename)
    pass
    


# # 包装

# In[8]:


def main(ppt_filename, output_filename, tts_engine):
    ve=init_tts(tts_engine)
    prs=read_pptx(ppt_filename)
    N_slides=len(prs.slides)
    
    with open(output_filename+".txt","w") as f:
        for index, slide in enumerate(prs.slides): 
            note=get_notes_text(slide)
            voice_filename=save_notes_voice(ve, note, index, tts_engine)
#             voice_filename=reclean_voice(voice_filename) # 清理一遍音频
            insert_voice(voice_filename, slide)
            print("Slide No. {} / {}".format(index+1,N_slides))
            clean_temp(voice_filename)
            f.write(note)
    prs.save(output_filename)
    print("save to ",output_filename)


# In[19]:


def interpret_opt(opt):
    # will rewrite with getopt
        
    tts_engine_dict={"Darwin":"nsss",
                     "Linux": "espeak",
                     "Windows":"sapi5",
                     "xunfei": "xunfei",
                     "google": "google"}
    
    if len(opt)==4:
        ppt_filename=opt[1]
        output_filename=opt[2]
        tts_engine_flag = opt[3]
        if opt[3]=="--online":
            tts_engine_flag="google"
#         else:
#             tts_engine_flag=platform.system()
        
    elif len(opt)==3 :
        ppt_filename=opt[1]
        output_filename=opt[2]
        tts_engine_flag=platform.system()
        if ppt_filename=="-f":
            ppt_filename="sample/test.pptx" # local test
            ppt_path=os.path.dirname(ppt_filename)
            output_filename=os.path.join(ppt_path, "output.pptx")




    elif len(opt)==2 :
        ppt_filename=opt[1]    
        ppt_path=os.path.dirname(ppt_filename)
        output_filename=os.path.join(ppt_path, "output.pptx")
        tts_engine_flag=platform.system()


    else:
        raise UserWarning("参数输入错误")
    tts_engine=tts_engine_dict[tts_engine_flag]


    print("Input file: ", ppt_filename)
    print("Output file: ", output_filename)
    print("TTS engine: ", tts_engine)
    return (ppt_filename, output_filename, tts_engine )


# In[20]:


if __name__=="__main__":
    ppt_filename, output_filename, tts_engine=interpret_opt(sys.argv)
    main(ppt_filename, output_filename, tts_engine)


# In[ ]:




