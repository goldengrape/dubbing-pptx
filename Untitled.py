
# coding: utf-8

# In[7]:


from google_tts import Speech as neo_Speech
from google_speech import Speech as ori_Speech
from io import BytesIO
from pydub import AudioSegment
from pydub.playback import play


text = '''
移除sox支持, 使之成为纯粹的python工具. 减少其他相关程序的安装. 受够了开源软件的一个又一个依赖. (已经完成)
'''
# This sentence has both Chinese and English.

lang = "zh-cn"
speech = neo_Speech(text,lang)
# speech.play()
# speech.save("test.mp3")
f=BytesIO()
speech.savef(f)


# configuration: 
# --prefix=/usr/local/Cellar/ffmpeg/4.1.1-with-options 
# --enable-shared 
# --enable-hardcoded-tables 
# --cc=clang 
# --host-cflags= --host-ldflags= 
# --enable-gpl 
# --enable-libaom 
# --enable-libmp3lame --enable-libopus --enable-libsnappy --enable-libtheora --enable-libvorbis --enable-libvpx --enable-libx264 --enable-libx265 --enable-libfontconfig --enable-libfreetype --enable-frei0r --enable-libass --disable-libjack --disable-indev=jack --enable-opencl --enable-videotoolbox --disable-htmlpages --enable-libfdk-aac --enable-libopenh264 --enable-librsvg --enable-libsrt --enable-libvidstab --enable-libxvid --enable-nonfree

# In[13]:


from pptx import Presentation
from pptx.util import Inches

ppt_filename="sample/test.pptx"
output_filename="sample/test_output.pptx"
prs = Presentation(ppt_filename)
left = top = Inches(0.0)
width = height = Inches(1.0)
for index, slide in enumerate(prs.slides): 
    shapes = slide.shapes
    movie = shapes.add_movie(f,  # 貌似不可以用
                             left , top , width , height, 
                             poster_frame_image=None, 
                             mime_type='video/unknown')
prs.save(output_filename)

